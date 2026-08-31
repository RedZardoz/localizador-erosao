import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_pptx_path):
    prs = Presentation()
    # 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette (Academic / Geospatial Dark Theme)
    BG_DARK = RGBColor(15, 23, 42)        # Slate 900
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800
    CARD_BORDER = RGBColor(51, 65, 85)    # Slate 700
    EMERALD_GREEN = RGBColor(16, 185, 129)# Emerald 500
    CYAN_ACCENT = RGBColor(6, 182, 212)   # Cyan 500
    INDIGO_ACCENT = RGBColor(99, 102, 241)# Indigo 500
    TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
    TEXT_LIGHT = RGBColor(203, 213, 225)  # Slate 300
    ACCENT_AMBER = RGBColor(245, 158, 11) # Amber 500
    ACCENT_ROSE = RGBColor(244, 63, 94)   # Rose 500
    ACCENT_BLUE = RGBColor(59, 130, 246)  # Blue 500

    def set_slide_background(slide, color=BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="PPGTCA 2026 • ESTRATÉGIA HÍBRIDA DE AMOSTRAGEM & IA"):
        # Category pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = EMERALD_GREEN

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title, content_lines, accent_color=EMERALD_GREEN, bg_color=CARD_BG, title_size=13, body_size=10.5):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.2)

        padding = Inches(0.2)
        tb = slide.shapes.add_textbox(left + padding, top + padding, width - (padding * 2), height - (padding * 2))
        tf = tb.text_frame
        tf.word_wrap = True

        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.size = Pt(title_size)
            p_title.font.bold = True
            p_title.font.color.rgb = accent_color
            p_title.space_after = Pt(6)

        first = not bool(title)
        for line in content_lines:
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.text = line
            p.font.size = Pt(body_size)
            p.font.color.rgb = TEXT_LIGHT
            p.space_after = Pt(4)

    def add_metric_badge(slide, left, top, width, height, label, value, unit="", color=EMERALD_GREEN):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = label.upper()
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MUTED

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = f"{value} {unit}".strip()
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = color

    # =============================================================
    # SLIDE 1: Capa
    # =============================================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    b_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.5))
    p_b = b_box.text_frame.paragraphs[0]
    p_b.text = "PROGRAMA DE PÓS-GRADUAÇÃO EM TECNOLOGIAS COMPUTACIONAIS PARA O AGRONEGÓCIO (PPGTCA - 2026)"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = EMERALD_GREEN

    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(2.2))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "Estratégia Híbrida de Amostragem e Rotulagem para Treinamento do Modelo Preditivo XGBoost"
    p_t.font.size = Pt(28)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE

    sub_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Integração entre Triagem GEE, Fotointerpretação em Alta Resolução (VHR) e Subamostra de Calibração Presencial (GNSS RTK / VANT)"
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = CYAN_ACCENT

    add_card(s1, Inches(1.0), Inches(5.4), Inches(11.333), Inches(1.4), "", [
        "Linha de Pesquisa: Sensoriamento Remoto, Inteligência Geoespacial e Predição de Erosão Laminar",
        "Objetivo Metodológico: Viabilizar dataset supervisionado estatisticamente robusto (N = 250 a 300 instâncias) com rigor científico auditável perante a banca de mestrado através de validação cruzada híbrida.",
    ], CYAN_ACCENT)

    # =============================================================
    # SLIDE 2: O Dilema Amostral
    # =============================================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "1. O Dilema Amostral no Aprendizado de Máquina Geoespacial")

    add_card(s2, Inches(0.8), Inches(1.6), Inches(3.6), Inches(5.2), "Abordagem 100% Campo", [
        "• Demanda: Visita presencial a 300 pontos com equipe e RTK.",
        "• Gargalo: Custo logístico e tempo de viagem inviáveis para o cronograma de mestrado.",
        "• Risco: Amostra insuficiente (N < 50) gerando subajuste (underfitting) no XGBoost.",
    ], ACCENT_ROSE)

    add_card(s2, Inches(4.8), Inches(1.6), Inches(3.6), Inches(5.2), "Abordagem 100% Sintética", [
        "• Demanda: Treinar o modelo direto nas fórmulas da RUSLE ou saídas do GEE.",
        "• Gargalo: Violação metodológica grave — o modelo apenas memoriza equações já conhecidas.",
        "• Risco: Rejeição pela banca por falta de dados empíricos reais de calibração.",
    ], ACCENT_AMBER)

    add_card(s2, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2), "Estratégia Híbrida Proposta", [
        "• Solução: Fotointerpretação VHR (30cm-3m) para 80% do dataset + Subamostra de Campo (20%).",
        "• Benefício: Volume ideal (N = 250-300) com custo operacional perfeitamente exequível.",
        "• Rigor: Matriz de Confusão e Índice Kappa para validação científica da rotulagem.",
    ], EMERALD_GREEN)

    # =============================================================
    # SLIDE 3: O Núcleo do Fluxo e Quantidades Amostrais
    # =============================================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "2. Arquitetura do Fluxo Amostral & Quantitativos por Etapa")

    add_metric_badge(s3, Inches(0.8), Inches(1.5), Inches(2.2), Inches(1.2), "Etapa 1: Triagem GEE", "250 a 300", "candidatos", INDIGO_ACCENT)
    add_metric_badge(s3, Inches(3.2), Inches(1.5), Inches(2.2), Inches(1.2), "Etapa 2: Fotointerp.", "100%", "(N = 250-300)", CYAN_ACCENT)
    add_metric_badge(s3, Inches(5.6), Inches(1.5), Inches(2.2), Inches(1.2), "Etapa 3: Campo RTK", "45", "pontos (20%)", EMERALD_GREEN)
    add_metric_badge(s3, Inches(8.0), Inches(1.5), Inches(2.2), Inches(1.2), "Etapa 4: Calibração", "Kappa ≥ 0.75", "(Acurácia ≥ 85%)", ACCENT_AMBER)
    add_metric_badge(s3, Inches(10.4), Inches(1.5), Inches(2.1), Inches(1.2), "Etapa 5: XGBoost", "250 a 300", "instâncias", ACCENT_BLUE)

    add_card(s3, Inches(0.8), Inches(2.9), Inches(11.7), Inches(3.9), "Tabela Consolidada do Fluxo Amostral", [
        "• ETAPA 1 (Triagem GEE): N = 250 a 300 candidatos gerados por máscara de elegibilidade + estratificação de relevo/solo + thinning de 500m.",
        "• ETAPA 2 (Fotointerpretação VHR): 100% dos candidatos inspecionados em imagens de satélite submétricas (Mapbox/Google 3D) com chave visual padronizada.",
        "• ETAPA 3 (Subamostra Presencial): n = 45 pontos (15% a 20%) sorteados de forma estratificada para visita física (GNSS RTK, VANT e KoboToolbox).",
        "• ETAPA 4 (Matriz de Concordância): Cálculo de Acurácia Global e Índice Kappa de Cohen entre fotointerpretação e medição real de campo.",
        "• ETAPA 5 (Treinamento do XGBoost): Particionamento 70% Treino (190 pts), 15% Validação (45 pts) e 15% Teste Cego Independente (45 pts).",
    ], CYAN_ACCENT, title_size=14, body_size=11)

    # =============================================================
    # SLIDE 4: Etapa 1 - Triagem GEE
    # =============================================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "3. Etapa 1: Triagem Estratificada no GEE (N = 250 a 300 Candidatos)")

    add_card(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Critérios Algorítmicos no Earth Engine", [
        "1. Máscara de Elegibilidade:",
        "   - Cobertura: ESA WorldCover v200 (10m) — focado em Cropland(40), Grassland(30) e Solo Exposto(60). Exclusão de áreas urbanas e floresta.",
        "   - Declividade: Copernicus DEM GLO-30 na faixa de 3% a 20% (classes de maior suscetibilidade laminar segundo Embrapa).",
        "   - Hidrografia: JRC Global Surface Water com buffer de exclusão.",
        "2. Estratificação Pedotopográfica:",
        "   - Cruzamento de 3 classes de declividade × 2 ordens de erodibilidade (SiBCS), formando 6 estratos representativos (A1 a B3).",
    ], INDIGO_ACCENT)

    add_card(s4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "Thinning Espacial e Saída Estruturada", [
        "3. Amostragem Estratificada Balanceada:",
        "   - Sorteio de ~50 pontos por estrato sobre a máscara de elegibilidade.",
        "4. Filtro de Thinning Espacial (d_min ≥ 500m):",
        "   - Algoritmo guloso de espaçamento euclidiano/haversine para evitar agrupamentos (clusters) e garantir dispersão homogênea na bacia.",
        "5. Metadados Extraídos no Request Único:",
        "   - Bandas Sentinel-2 (BSI, NDVI), elevação, declividade e fatores RUSLE aproximados anexados a cada candidato.",
        "• Resultado da Etapa: 250 a 300 pontos candidatos com status 'gee-screened'.",
    ], CYAN_ACCENT)

    # =============================================================
    # SLIDE 5: Etapa 2 - Fotointerpretação VHR
    # =============================================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "4. Etapa 2: Fotointerpretação em Alta Resolução VHR (100% dos Candidatos)")

    add_card(s5, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Chave Padronizada de Decisão Visual", [
        "• Classe 0 (Estável / Baixo Risco):",
        "  Solo com cobertura vegetal densa ou plantio direto consolidado com palhada abundante e terraços íntegros.",
        "• Classe 1 (Risco Moderado):",
        "  Solo exposto em preparo agrícola sobre relevo suave (< 6%), sem quebra de terraço ou manchas de fluxo.",
        "• Classe 2 (Erosão Laminar Severa):",
        "  Decapitação nítida do horizonte A (manchas avermelhadas/claras no topo da vertente), ausência ou rompimento de terraços.",
        "• Classe 3 (Erosão Crítica com Enxurrada):",
        "  Manchas claras de erosão convergindo para linhas de fluxo concentrado e acúmulo de sedimentos nas baixadas.",
    ], CYAN_ACCENT)

    add_card(s5, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "Fluxo Operacional e Rastreabilidade", [
        "• Ambiente de Inspeção:",
        "  Visualização direta no app em Zoom Z18-Z19 (Mapbox Satélite) e abertura em 1 clique no Google Earth Web 3D (pitch 45°).",
        "• Metadados Gravados para Cada Ponto:",
        "  - Evidências visuais (Solo exposto, decapitação, linhas de fluxo, falha de terraço, palhada).",
        "  - Grau de Confiança da Interpretação ('Alta' vs 'Média').",
        "  - Severidade Atribuída ('Moderada', 'Alta', 'Crítica').",
        "• Promoção de Status no Sistema:",
        "  O ponto passa de 'gee-screened' para 'vhr-photointerpreted'.",
        "• Amostras Processadas: 250 a 300 pontos rotulados.",
    ], EMERALD_GREEN)

    # =============================================================
    # SLIDE 6: Etapa 3 - Subamostra Presencial
    # =============================================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "5. Etapa 3: Subamostra Presencial de Calibração (n = 45 Pontos)")

    add_card(s6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Dimensionamento & Sorteio Estratificado", [
        "• Tamanho da Subamostra: n = 45 pontos (15% a 20% do total).",
        "• Critério de Sorteio:",
        "  - Distribuição proporcional entre os 6 estratos pedotopográficos.",
        "  - Inclusão balanceada das classes de severidade atribuídas na fotointerpretação.",
        "• Exportação de Missão de Campo:",
        "  - Geração de arquivo KoboToolbox (formulário digital de campo) e arquivo GPX/KML para navegação no receptor GNSS RTK e controladora do VANT.",
    ], EMERALD_GREEN)

    add_card(s6, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "Protocolo de Diagnóstico Físico Presencial", [
        "1. Pedestais de Erosão:",
        "   Medição milimétrica da altura do solo protegido sob pedregulhos/restos culturais (determina a lâmina decapitada).",
        "2. Descalçamento Radicular:",
        "   Avaliação de raízes expostas acima da superfície do solo.",
        "3. Selamento Superficial (Soil Crust):",
        "   Identificação de crosta compactada por impacto de gota.",
        "4. Registro Fotográfico Padronizado:",
        "   Fotos em ângulo nadir (1.5m) para estimativa de palhada e fotos panorâmicas do talhão.",
        "• Status no Sistema: Promovido a 'field-validated' (Ground-Truth).",
    ], ACCENT_AMBER)

    # =============================================================
    # SLIDE 7: Etapa 4 - Matriz de Confusão e Calibração
    # =============================================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "6. Etapa 4: Matriz de Confusão & Calibração Metodológica")

    add_card(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Cruzamento Fotointerpretação x Campo", [
        "• Pareamento dos 45 Pontos Visitados:",
        "  Comparação direta entre o rótulo atribuído por imagem de satélite VHR e o diagnóstico físico real medido em campo.",
        "• Matriz de Confusão 3x3:",
        "  Avaliação de acertos, falso-positivos e falso-negativos para as classes Moderada, Alta e Crítica.",
        "• Métricas Científicas de Validação:",
        "  - Acurácia Global (Meta: ≥ 85% de concordância).",
        "  - Índice Kappa de Cohen (Meta: κ ≥ 0.75 — concordância substancial/quase perfeita).",
        "  - Sensibilidade e Especificidade por classe de erosão.",
    ], ACCENT_AMBER)

    add_card(s7, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "O Valor desta Etapa para a Banca de Mestrado", [
        "• Resposta à Pergunta Central da Banca:",
        "  'Quão confiável é a rotulagem remota em alta resolução para treinar o algoritmo preditivo?'",
        "• Blindagem Metodológica:",
        "  Apresenta dados estatísticos concretos comprovando a correlação entre feições espectrais e processos físicos de campo.",
        "• Calibração de Limiares:",
        "  Caso o Kappa revele viés em alguma classe (e.g. subestimação em Latossolos), os rótulos restantes são re-calibrados antes do treino do XGBoost.",
    ], INDIGO_ACCENT)

    # =============================================================
    # SLIDE 8: Etapa 5 - XGBoost & SHAP
    # =============================================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "7. Etapa 5: Modelagem Preditiva XGBoost & Interpretabilidade SHAP")

    add_card(s8, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Particionamento e Treinamento do Modelo", [
        "• Volume Total do Dataset: N = 250 a 300 instâncias tabulares.",
        "• Divisão Estratificada:",
        "  - Treinamento (70% ≈ 190 instâncias): Ajuste dos pesos das árvores de decisão gradiente (Gradient Boosted Trees).",
        "  - Validação (15% ≈ 45 instâncias): Ajuste fino de hiperparâmetros (max_depth, learning_rate, subsample).",
        "  - Teste Cego (15% ≈ 45 instâncias): Avaliação final com amostras reservadas de campo (Ground-Truth).",
        "• Variáveis de Entrada (X):",
        "  BSI, NDVI, Declividade (%), Elevação (m), Fator K (solo), Fator LS (relevo), Fator R (chuva), Fator C (manejo).",
    ], ACCENT_BLUE)

    add_card(s8, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "Interpretabilidade e Diagnóstico via SHAP", [
        "• Algoritmo SHAP (SHapley Additive exPlanations):",
        "  Calcula a contribuição marginal de cada fator biofísico na predição de risco de erosão para cada talhão.",
        "• Entregáveis para a Dissertação:",
        "  - Gráfico SHAP Summary Plot: Ranking de importância das variáveis no território paranaense.",
        "  - Gráficos de Dependência Parcial: Limiares críticos de declividade e BSI onde a perda de solo dispara.",
        "  - Métricas de Desempenho: Acurácia, F1-Score, Curva ROC-AUC e Erro Médio Quadrático (RMSE).",
    ], EMERALD_GREEN)

    # =============================================================
    # SLIDE 9: Recursos no Aplicativo
    # =============================================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "8. Implementação Técnica no Aplicativo Web")

    add_card(s9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Novos Componentes e Estruturas de Dados", [
        "1. Tipagem e Proveniência (src/types/erosion.ts):",
        "   Novo status 'vhr-photointerpreted' e interface PhotoInterpretationRecord com evidências visuais.",
        "2. Gaveta de Fotointerpretação Assistida:",
        "   Formulário rápido acoplado ao popup do ponto com atalhos de classificação, zoom automático e integração Google Earth 3D.",
        "3. Seletor da Subamostra de Campo:",
        "   Sorteio estratificado dos 45 pontos com exportação direta para KoboToolbox / GPX.",
    ], CYAN_ACCENT)

    add_card(s9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), "Módulo de Calibração e Exportador", [
        "4. Painel de Matriz de Confusão (KoboFieldImport.tsx):",
        "   Cálculo automático de Acurácia Global e Índice Kappa na importação dos dados de campo.",
        "5. Exportação Dedicada para IA (ExportModal.tsx):",
        "   Geração do arquivo 'dataset_treinamento_xgboost.csv' consolidando variáveis biofísicas, fatores RUSLE e coluna de origem do rótulo.",
        "6. Testes Automatizados (Vitest):",
        "   Cobertura de testes para garantir integridade do cálculo de concordância e transições de status.",
    ], INDIGO_ACCENT)

    # =============================================================
    # SLIDE 10: Conclusão e Cronograma
    # =============================================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, "9. Síntese Amostral e Conclusão Metodológica")

    add_card(s10, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2), "Resumo Executivo da Estratégia Híbrida", [
        "• Redução de Custo & Tempo: Reduz a necessidade de deslocamento em 80% mantendo representatividade espacial.",
        "• Rigor Metodológico Auditável: A subamostra de 45 pontos com RTK/Drone e a Matriz de Confusão fornecem a comprovação científica necessária para a defesa da dissertação de mestrado.",
        "• Volume Amostral Adequado: N = 250 a 300 instâncias garantem convergência estável e generalização confiável para o XGBoost.",
        "• Cronograma Sugerido: Etapa 1 e 2 (Semanas 1-2) | Etapa 3 (Semanas 3-4) | Etapas 4 e 5 (Semanas 5-6).",
        "• Status do Projeto: Pronto para implantação no código da aplicação web.",
    ], EMERALD_GREEN, title_size=15, body_size=12)

    prs.save(output_pptx_path)
    print(f"Presentation saved successfully to: {output_pptx_path}")

if __name__ == "__main__":
    out_path = os.path.join(os.getcwd(), "Plano_Estrategia_Hibrida_XGBoost_PPGTCA.pptx")
    create_presentation(out_path)
