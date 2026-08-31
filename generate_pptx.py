import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_pptx_path):
    prs = Presentation()
    # 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette (Academic / Geospatial Dark Theme)
    BG_DARK = RGBColor(15, 23, 42)        # Slate 900
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800
    CARD_BORDER = RGBColor(51, 65, 85)    # Slate 700
    EMERALD_GREEN = RGBColor(16, 185, 129)# Emerald 500
    CYAN_ACCENT = RGBColor(6, 182, 212)   # Cyan 500
    INDIGO_ACCENT = RGBColor(99, 102, 241)# Indigo 500
    TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
    TEXT_LIGHT = RGBColor(203, 213, 225)  # Slate 300
    ACCENT_AMBER = RGBColor(245, 158, 11) # Amber 500
    ACCENT_ROSE = RGBColor(244, 63, 94)   # Rose 500

    def set_slide_background(slide, color=BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="PPGTCA 2026 • METODOLOGIA CIENTÍFICA"):
        # Category pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = EMERALD_GREEN

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title, content_lines, accent_color=EMERALD_GREEN):
        # Shape background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.2)

        # Content text
        padding = Inches(0.25)
        tb = slide.shapes.add_textbox(left + padding, top + padding, width - (padding * 2), height - (padding * 2))
        tf = tb.text_frame
        tf.word_wrap = True

        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.size = Pt(14)
            p_title.font.bold = True
            p_title.font.color.rgb = accent_color
            p_title.space_after = Pt(8)

        first = not bool(title)
        for line in content_lines:
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT
            p.space_after = Pt(5)

    # -------------------------------------------------------------
    # SLIDE 1: Capa
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Badge
    b_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.5))
    p_b = b_box.text_frame.paragraphs[0]
    p_b.text = "PROGRAMA DE PÓS-GRADUAÇÃO EM TECNOLOGIAS COMPUTACIONAIS PARA O AGRONEGÓCIO (PPGTCA - 2026)"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = EMERALD_GREEN

    # Title
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "Metodologia Geoespacial de Detecção, Triagem e Priorização de Focos de Erosão Laminar"
    p_t.font.size = Pt(30)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE

    # Subtitle
    sub_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Pipeline Algorítmico no Google Earth Engine, Estratificação Físico-Pedológica e Amostragem para Coletas de Campo no Paraná"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = CYAN_ACCENT

    # Footer Card
    add_card(s1, Inches(1.0), Inches(5.4), Inches(11.333), Inches(1.3), "", [
        "Linha de Pesquisa: Sensoriamento Remoto, Inteligência Geoespacial e Conservação de Recursos Naturais",
        "Objetivo da Apresentação: Detalhamento minucioso do fluxo metodológico e procedimentos algorítmicos no Earth Engine para averiguação e orientação científica.",
    ], CYAN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 2: Problema Científico e Motivação
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "1. Contextualização Científica & O Desafio da Erosão Laminar")

    add_card(s2, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "O Fenômeno da Erosão Laminar (Sheet Erosion)", [
        "• Processo Contínuo e Difuso: Desprendimento uniforme das partículas do solo pelo impacto das gotas de chuva (splash detachment) e arraste por escoamento superficial não-concentrado.",
        "• Ausência de Incisões Iniciais: Diferente de ravinas e voçorocas, não apresenta cicatrizes visíveis imediatas no relevo, dificultando detecção prévia.",
        "• Impacto Agronômico Silencioso: Degradação da camada mais fértil (horizonte A), perda de carbono orgânico e redução da capacidade de retenção hídrica.",
        "• Custo Regional: Contaminação e assoreamento de mananciais nas bacias hidrográficas do Paraná (Iguaçu, Tibagi, Ivaí, Paraná 3).",
    ], ACCENT_AMBER)

    add_card(s2, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "A Proposta Metodológica da Pesquisa", [
        "• Inviabilidade de Amostragem Aleatória: O estado do Paraná possui 199.315 km²; coletar dados de campo às cegas gera custo proibitivo e baixa representatividade.",
        "• Triagem Geoespacial Determinística: Desenvolver um pipeline no Google Earth Engine que selecione talhões-piloto candidatos com rigor estatístico e físico.",
        "• Base para Inteligência Artificial: Os pontos pré-triados servem como alvos preliminares para validação em campo (GNSS RTK / VANT), gerando o dataset rotulado para o classificador XGBoost/SHAP.",
        "• Auditabilidade Total: Cada etapa (máscara, estratificação, thinning e RUSLE) possui rastreabilidade matemática aberta.",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 3: Arquitetura Geral do Pipeline
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "2. Arquitetura Geral do Pipeline Metodológico de Seleção")

    col_w = Inches(2.2)
    gap = Inches(0.18)
    top_pos = Inches(1.7)
    h_pos = Inches(5.0)

    add_card(s3, Inches(0.8) + (col_w + gap) * 0, top_pos, col_w, h_pos, "Etapa 1: AOI", [
        "Delimitação Espacial",
        "• Polígono IBGE oficial (Município / Microbacia)",
        "• Upload de AOI GeoJSON/KML",
        "• Gradeamento automático em tiles caso AOI > 0.5°",
    ], CYAN_ACCENT)

    add_card(s3, Inches(0.8) + (col_w + gap) * 1, top_pos, col_w, h_pos, "Etapa 2: Elegibilidade", [
        "Máscara 10m (GEE)",
        "• Uso da Terra (WorldCover)",
        "• Declividade 3% a 20% (Copernicus DEM métrico)",
        "• Buffer de 30m de corpos hídricos (JRC Water)",
    ], EMERALD_GREEN)

    add_card(s3, Inches(0.8) + (col_w + gap) * 2, top_pos, col_w, h_pos, "Etapa 3: Estrato", [
        "Matriz Cruzada 2x3",
        "• Relevo (≤6%, 6-12%, >12%)",
        "• Pedologia (K Grupo A / B)",
        "• Amostragem balanceada de sub-estratos (A1 a B3)",
    ], INDIGO_ACCENT)

    add_card(s3, Inches(0.8) + (col_w + gap) * 3, top_pos, col_w, h_pos, "Etapa 4: Thinning", [
        "Desaglomeração Espacial",
        "• Distância geodésica mínima de Haversine (ex: 1.0 km)",
        "• Preserva focos com maior Priority Score",
        "• Elimina autocorrelação",
    ], ACCENT_AMBER)

    add_card(s3, Inches(0.8) + (col_w + gap) * 4, top_pos, col_w, h_pos, "Etapa 5: Validação", [
        "RUSLE & Campo",
        "• Fatores R, K, LS, C reais",
        "• Importação KoboToolbox",
        "• Exportação de Dataset Rotulado para IA",
    ], ACCENT_ROSE)

    # -------------------------------------------------------------
    # SLIDE 4: Fase 1 - Máscara de Uso da Terra (ESA WorldCover)
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "3. Máscara de Elegibilidade — Critério 1: Uso da Terra (ESA WorldCover 10m)")

    add_card(s4, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Base de Dados e Resolução", [
        "• Coleção GEE: 'ESA/WorldCover/v200/2021'",
        "• Resolução Espacial: 10 metros",
        "• Sensor Base: Fusão multiespectral Sentinel-1 SAR + Sentinel-2 MSI",
        "• Classes Elegíveis para Triagem:",
        "  - Cropland (40): Lavouras anuais e perenes sob manejo agrícola;",
        "  - Grassland (30): Pastagens cultivadas e campos limpos;",
        "  - Bare / sparse vegetation (60): Solo exposto e vegetação esparsa.",
    ], CYAN_ACCENT)

    add_card(s4, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Justificativa Científica & Exclusões", [
        "• Classes Excluídas Rigorosamente:",
        "  - Tree cover (10): Florestas densas e matas ciliares (alta interceptação e rugosidade superficial);",
        "  - Built-up (50): Superfícies impermeáveis urbanas (outro mecanismo de escoamento);",
        "  - Permanent water bodies (80): Corpos d'água permanentes.",
        "• Algoritmo GEE de Filtragem:",
        "  - Remapeamento binário com ee.Image.remap()",
        "  - lcMask = remap([30, 40, 60], [1, 1, 1], 0)",
        "  - Garante foco restrito a superfícies com potencial real de erosão hídrica em áreas rurais.",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 5: Fase 1 - Máscara de Declividade (Copernicus DEM)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "4. Máscara de Elegibilidade — Critério 2: Declividade (Copernicus DEM GLO-30)")

    add_card(s5, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Equação e Modelagem Físico-Espacial", [
        "• Coleção: 'COPERNICUS/DEM/GLO30' (Resolução 30m)",
        "• Ancoragem Geométrica Métrica: Aplicação de .setDefaultProjection('EPSG:3857', null, 30) para garantir que Δz e Δx estejam em metros.",
        "• Cálculo da Declividade em Graus:",
        "  θ = ee.Terrain.slope(DEM_metrico)",
        "• Conversão Trigonométrica para Porcentagem (S%):",
        "  S% = tan(θ · π / 180) × 100",
        "• Faixa de Elegibilidade Paramétrica:",
        "  3.0% ≤ S% ≤ 20.0%",
    ], EMERALD_GREEN)

    add_card(s5, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Fundamentação Agronômica e Pedológica", [
        "• Por que excluir declividades < 3%?",
        "  Áreas extremamente planas têm predomínio de escoamento muito lento e deposição/infiltração, sem energia cinética suficiente para arraste difuso severo.",
        "• Por que limitar em 20%?",
        "  Declividades > 20% no SiBCS (relevo forte-ondulado a montanhoso) favorecem rápida concentração em sulcos profundos e ravinamento, saindo do escopo estrito de erosão laminar.",
        "• Referência Normativa:",
        "  Bertoni & Lombardi Neto (2017) e Critérios de Aptidão Agrícola das Terras (Embrapa/SiBCS).",
    ], ACCENT_AMBER)

    # -------------------------------------------------------------
    # SLIDE 6: Fase 1 - Exclusão Hídrica (JRC Global Surface Water)
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "5. Máscara de Elegibilidade — Critério 3: Corpos d'Água e Buffer de 30m")

    add_card(s6, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Dataset e Extração Hidrográfica", [
        "• Coleção GEE: 'JRC/GSW1_4/GlobalSurfaceWater'",
        "• Banda de Análise: 'occurrence' (0 a 100% de persistência hídrica multitemporal histórica).",
        "• Limiar de Identificação de Água: occurrence > 10%",
        "• Aplicação de Buffer Morfológico de Segurança:",
        "  waterMask = occurrence.gt(10)",
        "  waterBuffered = waterMask.focalMax(30, 'circle', 'meters')",
        "• waterEligible = waterBuffered.unmask(0).not()",
    ], CYAN_ACCENT)

    add_card(s6, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Justificativa e Prevenção de Falsos Positivos", [
        "• Eliminação de Assinaturas Espectrais Úmidas: Áreas ribeirinhas, várzeas e banhados possuem alta umidade e matéria orgânica que podem falsear índices espectrais de solo.",
        "• Proteção de Áreas de Preservação Permanente (APP): Respeito ao Código Florestal Brasileiro (Lei 12.651/2012) — faixas marginais de cursos d'água naturais.",
        "• Isolamento Hidrológico: O buffer de 30m garante que os candidatos amostrados reflitam vertentes agrícolas de lavoura e não erosão de margem fluvial.",
    ], INDIGO_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 7: Interseção e Diagnóstico Passo 0
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "6. Interseção Lógica da Máscara & Diagnóstico de Pixels (Passo 0)")

    add_card(s7, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Interseção Booleana Raster (10m)", [
        "• Equação de Elegibilidade Total:",
        "  Eligibility = lcMask ∧ slopeEligible ∧ waterEligible",
        "• Recorte Espacial Vetorial:",
        "  EligibilityMasked = Eligibility.updateMask(Eligibility.eq(1)).clip(AOI)",
        "• Aplicação da Máscara no Grafo de Computação:",
        "  Todas as bandas biofísicas (BSI, NDVI, DEM, Slope, Estrato) herdam a máscara ativa antes da amostragem.",
    ], EMERALD_GREEN)

    add_card(s7, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Instrumentação Diagnóstica Empírica", [
        "• Reducer de Auditoria no Servidor (reduceRegion a 30m):",
        "  - lc_sum: Pixels que atendem ao uso do solo;",
        "  - slope_sum: Pixels que atendem à declividade;",
        "  - water_sum: Pixels fora de corpos d'água e buffer;",
        "  - eligible_sum: Pixels onde os 3 critérios são simultaneamente válidos.",
        "• Evidência Empírica em Céu Azul (PR):",
        "  - Uso do Solo: 288.829 pixels",
        "  - Fora d'Água: 1.441.895 pixels",
        "  - Diagnóstico transparente exposto ao pesquisador na interface web.",
    ], CYAN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 8: Fase 2 - Estratificação Cruzada
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "7. Fase 2 — Estratificação Cruzada Amostral (Matriz 2 x 3)")

    add_card(s8, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "A Matriz Estruturada de 6 Sub-Estratos", [
        "• Eixo 1: Classes de Relevo (Topografia):",
        "  - Classe 1: Suave (Declividade ≤ 6%)",
        "  - Classe 2: Moderada (6% < Declividade ≤ 12%)",
        "  - Classe 3: Forte (Declividade > 12%)",
        "• Eixo 2: Grupos Pedológicos de Erodibilidade (K):",
        "  - Grupo A (Alta Erodibilidade): Argissolos Vermelho-Amarelos, Neossolos Litólicos/Regolíticos (K ≈ 0.045)",
        "  - Grupo B (Média/Baixa Erodibilidade): Latossolos Vermelhos, Nitossolos (K ≈ 0.020)",
        "• Banda Raster GEE: 'STRATUM' com valores 1 a 6 (A1, A2, A3, B1, B2, B3).",
    ], INDIGO_ACCENT)

    add_card(s8, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Amostragem Estratificada no GEE", [
        "• Função no Earth Engine:",
        "  image.stratifiedSample({",
        "    numPoints: ceil(targetCount / 6) * 3,",
        "    classBand: 'STRATUM',",
        "    region: AOI,",
        "    scale: 30,",
        "    dropNulls: true,",
        "    geometries: true",
        "  })",
        "• Vantagem Científica: Evita o viés de superamostragem em áreas homogêneas e garante representatividade de todos os tipos de terreno.",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 9: Fase 3 - Thinning Espacial Geodésico
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "8. Fase 3 — Thinning Espacial Geodésico (Desaglomeração)")

    add_card(s9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "O Problema da Autocorrelação Espacial", [
        "• Primeira Lei da Geografia (Tobler): Pixels vizinhos tendem a ter valores de solo, declividade e BSI muito semelhantes.",
        "• Se a amostragem coletar 5 pontos no mesmo talhão agrícola a 50 metros um do outro, a coleta de campo terá redundância e desperdício logístico.",
        "• Solução: Algoritmo guloso de Thinning Espacial por Distância Mínima.",
    ], ACCENT_AMBER)

    add_card(s9, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Algoritmo de Haversine & Priority Score", [
        "1. Ordenação Decrescente: Todos os candidatos amostrados são ordenados pelo PriorityScore (0 a 100).",
        "2. Varredura Iterativa: Para cada candidato Pi:",
        "   - Calcula a distância geodésica dj para todos os pontos já aceitos;",
        "   - Se dj ≥ minSpacingKm (ex: 1.0 km) para todo j, Pi é mantido;",
        "   - Caso contrário, Pi é descartado.",
        "3. Resultado: Cobertura territorial ampla, dispersa e focada nos pontos mais críticos de cada microbacia.",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 10: Modelagem RUSLE e Fator R (NASA POWER)
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, "9. Modelagem Físico-Matemática — Equação RUSLE e Fator R")

    add_card(s10, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Equação Universal de Perda de Solo (RUSLE)", [
        "$$A = R \\times K \\times LS \\times C \\times P$$",
        "• A: Perda média anual de solo (t/ha/ano);",
        "• R: Erosividade da chuva (MJ·mm / ha·h·ano);",
        "• K: Erodibilidade do solo (t·h / ha·MJ·mm);",
        "• LS: Fator topográfico comprimento-declividade;",
        "• C: Fator de cobertura vegetal e manejo;",
        "• P: Fator de práticas conservacionistas de suporte.",
    ], CYAN_ACCENT)

    add_card(s10, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Fator R — Climatologia Real NASA POWER", [
        "• Fonte de Dados: Reanálise MERRA-2 (NASA POWER API) série temporal 2001-2020.",
        "• Equação de Lombardi Neto & Moldenhauer (1992) calibrada para o Sul/Sudeste:",
        "  R = 68.730 × (p² / P)^0.841",
        "• p = Precipitação média do mês mais chuvoso (mm);",
        "• P = Precipitação média anual acumulada (mm);",
        "• Substitui estimativas estáticas por climatologia local real e sem custos de API.",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 11: Índices Espectrais e Fator C
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s11)
    add_header(s11, "10. Índices Espectrais de Satélite (Sentinel-2) e Fator C")

    add_card(s11, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Bare Soil Index (BSI) & NDVI", [
        "• Cena: Copernicus Sentinel-2 MSI Harmonized L2A com máscara de nuvens SCL.",
        "• Bare Soil Index (BSI):",
        "  BSI = [(B12 + B4) - (B8 + B2)] / [(B12 + B4) + (B8 + B2)]",
        "• BSI > +0.35: Solo totalmente exposto e pulverizado (máxima suscetibilidade);",
        "• BSI entre +0.05 e +0.35: Cobertura esparsa ou palhada residual;",
        "• NDVI = (B8 - B4) / (B8 + B4): Vigor fotossintético.",
    ], CYAN_ACCENT)

    add_card(s11, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Derivação do Fator C & Severidade", [
        "• Fator C derivado espectralmente:",
        "  - NDVI > 0.6: C ≈ 0.02 (cobertura densa)",
        "  - NDVI entre 0.3 e 0.6: C = 0.15 a 0.35",
        "  - NDVI < 0.3 e BSI > 0.2: C ≈ 0.60 a 0.85 (solo desnudo)",
        "• Índice de Severidade:",
        "  Cruza Declividade (S%), BSI e Erodibilidade (K) para classificar em Moderada, Alta ou Crítica.",
        "• Score de Prioridade (0-100): Pondera Severidade (50%), BSI (25%), Declividade (25%).",
    ], INDIGO_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 12: Arquitetura de Software e Segurança
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s12)
    add_header(s12, "11. Arquitetura da Aplicação Web e Segurança de Credenciais")

    add_card(s12, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Stack Tecnológica e Frontend 3D", [
        "• Framework: Next.js 14 (App Router) + React 18 + TypeScript estrito.",
        "• Motor Cartográfico: MapLibre GL JS v4+ acelerado por WebGL.",
        "• Terreno Tridimensional: Malha DEM Terrarium AWS com elevação e iluminação azimutal em tempo real.",
        "• Gerenciamento de Estado: Zustand com seletores memoizados para performance em 60 FPS.",
    ], CYAN_ACCENT)

    add_card(s12, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Segurança Criptográfica de Sessão", [
        "• Autenticação Google OAuth2 Real:",
        "  Assinatura de JWT RS256 via módulo nativo crypto do Node.js trocado por token OAuth2.",
        "• Sessão Temporária de Servidor:",
        "  A chave privada RSA da Service Account NUNCA é salva em localStorage. Ela vive apenas na memória do servidor associada a um cookie httpOnly + secure + sameSite=strict.",
        "• Zero vazamento em repositórios ou logs.",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 13: Validação em Campo (KoboToolbox)
    # -------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s13)
    add_header(s13, "12. Ciclo de Validação em Campo — KoboToolbox → Dataset Rotulado")

    add_card(s13, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Protocolo de Coleta em Campo", [
        "• Dispositivos: GNSS Geodésico RTK (precisão centimétrica) e VANT Multiespectral.",
        "• Formulário Digital: Coleta estruturada no KoboToolbox / ODK Collect.",
        "• Indicadores de Campo Registrados:",
        "  - Altura de pedestais de erosão (mm);",
        "  - Descalçamento radicular de plântulas;",
        "  - Selamento superficial (crosta de solo);",
        "  - Espessura de lâminas de deposição.",
    ], ACCENT_AMBER)

    add_card(s13, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Integração e Pareamento Automático", [
        "• Parser Inteligente (koboParser.ts):",
        "  Lê o export CSV do KoboToolbox e casa cada registro com o foco pré-triado mais próximo (< 150m) ou por código de talhão.",
        "• Atualização de Proveniência: O ponto passa de 'gee-screened' para 'field-validated'.",
        "• Exportação para IA: Gera o dataset de treino/teste rotulado para os futuros modelos de Machine Learning (XGBoost e SHAP).",
    ], EMERALD_GREEN)

    # -------------------------------------------------------------
    # SLIDE 14: Sistema de Coleções Salvas
    # -------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s14)
    add_header(s14, "13. Sistema de Projetos & Coleções Salvas no Navegador")

    add_card(s14, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Persistência e Gerenciamento Local", [
        "• Histórico de Campanhas:",
        "  O pesquisador pode salvar a seleção de candidatos no navegador com nome de campanha, data e notas metodológicas.",
        "• Recarregamento Instantâneo com 1 Clique:",
        "  Restaura os pontos, a AOI delimitada e os filtros, acionando o enquadramento suave da câmera 3D.",
        "• Portabilidade Total:",
        "  Exportação e importação de arquivos de projeto .json entre diferentes computadores de pesquisa.",
    ], CYAN_ACCENT)

    add_card(s14, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Interoperabilidade Científica Multiformato", [
        "• Formatos Exportáveis:",
        "  - GeoJSON: Camadas com tabela de atributos para SIGs desktop (QGIS, ArcGIS Pro);",
        "  - KML 3D: Visualização imersiva no Google Earth Pro com popups formatados;",
        "  - CSV Científico: Coordenadas em DD e DMS, fatores RUSLE e variáveis para R/Python;",
        "  - Dataset XGBoost: Apenas pontos validados em campo para modelagem preditiva.",
    ], INDIGO_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 15: Rigor Metodológico e Suíte de Testes
    # -------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s15)
    add_header(s15, "14. Rigor Metodológico, Auditoria e Suíte de Testes")

    add_card(s15, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Testes Automatizados (Vitest)", [
        "• Cobertura Total do Pipeline:",
        "  14 arquivos de teste com 99 testes unitários aprovados em ~4.3s.",
        "• Módulos Testados Rigorosamente:",
        "  - Fórmulas RUSLE (rusleCalculator.test.ts);",
        "  - Autenticação e assinatura RSA JWT (googleAuth.test.ts);",
        "  - Thinning espacial de Haversine (spatialThinning.test.ts);",
        "  - Estratificação cruzada (stratification.test.ts);",
        "  - Pareamento de campo Kobo (koboParser.test.ts);",
        "  - Persistência e restauração de coleções (useErosionStore.test.ts).",
    ], EMERALD_GREEN)

    add_card(s15, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Garantias de Confiabilidade", [
        "• Tipagem Estrita TypeScript (npx tsc --noEmit: 0 erros).",
        "• Build de Produção Validado (Next.js 14 estático + dinâmico compilado sem falhas).",
        "• Reprodutibilidade: Qualquer pesquisador pode reproduzir os passos do GEE e obter exatamente os mesmos candidatos.",
    ], CYAN_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 16: Conclusões e Próximos Passos
    # -------------------------------------------------------------
    s16 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s16)
    add_header(s16, "15. Conclusões e Próximos Passos no PPGTCA")

    add_card(s16, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), "Contribuições Científicas Alcançadas", [
        "1. Pipeline Determinístico e Automatizado: Substituição de triagens manuais por um fluxo GEE estruturado.",
        "2. Amostragem Balanceada: Estratificação cruzada A1..B3 que evita viés espacial e geológico.",
        "3. Otimização Logística: Desaglomeração por thinning que maximiza a cobertura das saídas de campo.",
        "4. Conformidade Científica: Rastreabilidade aberta dos fatores biofísicos da RUSLE.",
    ], EMERALD_GREEN)

    add_card(s16, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), "Cronograma dos Próximos Passos", [
        "• Etapa 1: Execução das coletas de campo nas áreas pré-triadas (Céu Azul, Paranavaí) com formulário KoboToolbox.",
        "• Etapa 2: Ingestão dos dados de campo e fechamento do dataset rotulado.",
        "• Etapa 3: Treinamento e validação dos modelos de Machine Learning (XGBoost / Random Forest).",
        "• Etapa 4: Interpretação global e local dos fatores preditivos via valores SHAP (*SHapley Additive exPlanations*).",
    ], CYAN_ACCENT)

    prs.save(output_pptx_path)
    print(f"Presentation saved successfully to: {output_pptx_path}")

if __name__ == "__main__":
    output_path = os.path.abspath("Apresentacao_Metodologia_Erosao_PPGTCA.pptx")
    create_presentation(output_path)
